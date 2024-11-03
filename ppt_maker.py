# Create a new presentation with references from Langmead et al. (2018) and Wang et al. (2020)

from pptx import Presentation

# Create a new presentation
prs_new = Presentation()

# Slide 1 - Title Slide
slide_layout = prs_new.slide_layouts[0]
slide = prs_new.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Cloud-based and Distributed Computing Platforms for Genomic Data Handling"
subtitle = slide.placeholders[1]
subtitle.text = "Using cloud infrastructure for collaborative research and large dataset analysis"

# Slide 2 - Introduction to Genomics and Biological Data (Langmead & Nellore, 2018)
slide_layout = prs_new.slide_layouts[1]
slide = prs_new.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Introduction to Genomics and Biological Data"
content = slide.shapes.placeholders[1].text_frame
content.text = "• Genomic data is complex, requiring vast computational resources to analyze."
content.add_paragraph().text = "• Langmead & Nellore (2018) describe the increase in genomic data volume and complexity, highlighting cloud computing as a solution for processing and storage."

# Slide 3 - Cloud-based Platforms (Wang et al. 2020)
slide = prs_new.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Cloud-based and Distributed Computing Platforms"
content = slide.shapes.placeholders[1].text_frame
content.text = "• Cloud platforms offer scalability for storing and analyzing genomic data."
content.add_paragraph().text = "• Wang et al. (2020) discuss the role of cloud platforms like AWS and Google Cloud in genomic data analysis and global collaboration."

# Slide 4 - Case Study: Genomics Research (Langmead & Nellore, 2018)
slide = prs_new.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Case Study: Genomics Research in the Cloud"
content = slide.shapes.placeholders[1].text_frame
content.text = "• Genomics research projects like 1000 Genomes have utilized cloud platforms."
content.add_paragraph().text = "• According to Langmead & Nellore (2018), cloud computing enabled seamless collaboration across multiple research centers."

# Slide 5 - Technical Implementation (Wang et al. 2020)
slide = prs_new.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Technical Implementation"
content = slide.shapes.placeholders[1].text_frame
content.text = "• Tools such as Apache Spark and GATK are used to handle large-scale genomic datasets in the cloud."
content.add_paragraph().text = "• Wang et al. (2020) emphasize the importance of distributed computing systems like Hadoop for processing parallel genomic workloads."

# Slide 6 - Impact and Results (Langmead & Nellore, 2018)
slide = prs_new.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Impact and Results"
content = slide.shapes.placeholders[1].text_frame
content.text = "• Cloud platforms have accelerated genomic research and increased collaboration."
content.add_paragraph().text = "• Langmead & Nellore (2018) state that cloud solutions have improved the speed of research and made real-time collaboration more feasible."

# Slide 7 - Conclusion (Wang et al. 2020 and Langmead & Nellore, 2018)
slide = prs_new.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Conclusion"
content = slide.shapes.placeholders[1].text_frame
content.text = "• Cloud computing and distributed platforms are crucial for handling large-scale genomic data."
content.add_paragraph().text = "• Both Langmead & Nellore (2018) and Wang et al. (2020) highlight how these platforms will continue to play an essential role in genomic research."

# Save the updated presentation
pptx_file_new = "Cloud_Computing_Genomics_Presentation_Updated_References.pptx"
prs_new.save(pptx_file_new)
